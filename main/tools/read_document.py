import re
from pathlib import Path

# Импортируем функцию поиска файлов из file_operations (без дублирования кода)
try:
    from main.commands.file_operations import find_file
except ImportError:
    from commands.file_operations import find_file

# Опциональные библиотеки для чтения документов
try:
    from docx import Document as DocxDocument
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    import PyPDF2
    HAS_PDF = True
except ImportError:
    try:
        import pypdf as PyPDF2
        HAS_PDF = True
    except ImportError:
        HAS_PDF = False

try:
    import openpyxl
    HAS_XLSX = True
except ImportError:
    HAS_XLSX = False

try:
    from pptx import Presentation as PptxPresentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

# Максимальная длина текста для возврата
MAX_TEXT_LENGTH = 8000


def _read_txt(file_path: Path) -> str:
    """Читает текстовый файл."""
    encodings = ['utf-8', 'cp1251', 'cp866', 'latin-1']
    for enc in encodings:
        try:
            return file_path.read_text(encoding=enc)
        except (UnicodeDecodeError, UnicodeError):
            continue
    raise ValueError(f"Не удалось прочитать файл в известных кодировках")


def _read_docx(file_path: Path) -> str:
    """Читает .docx файл."""
    if not HAS_DOCX:
        return "[Ошибка: для чтения .docx установите python-docx]"
    
    doc = DocxDocument(str(file_path))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    return "\n".join(paragraphs)


def _read_doc(file_path: Path) -> str:
    """Читает .doc файл (старый формат Word)."""
    # Попробуем через antiword (если установлен в системе)
    try:
        import subprocess
        result = subprocess.run(
            ['antiword', str(file_path)],
            capture_output=True,
            text=True,
            timeout=10
        )
        if result.returncode == 0:
            return result.stdout
    except Exception:
        pass
    
    # Альтернатива: попытка прочитать как текст (иногда работает)
    try:
        content = file_path.read_bytes()
        # Извлекаем читаемый текст из бинарника
        text = content.decode('cp1251', errors='ignore')
        # Убираем нечитаемые символы
        text = re.sub(r'[^\x20-\x7E\u0400-\u04FF\n\r\t]', ' ', text)
        text = re.sub(r'\s+', ' ', text).strip()
        if len(text) > 100:  # Если удалось извлечь что-то осмысленное
            return text
    except Exception:
        pass
    
    return "[Ошибка: для чтения .doc нужен antiword или конвертация в .docx]"


def _read_pdf(file_path: Path) -> str:
    """Читает .pdf файл."""
    if not HAS_PDF:
        return "[Ошибка: для чтения .pdf установите PyPDF2 или pypdf]"
    
    text_parts = []
    with open(file_path, 'rb') as f:
        reader = PyPDF2.PdfReader(f)
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                text_parts.append(page_text)
    
    return "\n".join(text_parts)


def _read_xlsx(file_path: Path) -> str:
    """Читает .xlsx файл (Excel)."""
    if not HAS_XLSX:
        return "[Ошибка: для чтения .xlsx установите openpyxl]"
    
    wb = openpyxl.load_workbook(str(file_path), read_only=True, data_only=True)
    parts = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(cells):
                rows.append("\t".join(cells))
        if rows:
            parts.append(f"--- Лист: {sheet_name} ---\n" + "\n".join(rows))
    wb.close()
    return "\n\n".join(parts) if parts else "(Пустая таблица)"


def _read_pptx(file_path: Path) -> str:
    """Читает .pptx файл (PowerPoint)."""
    if not HAS_PPTX:
        return "[Ошибка: для чтения .pptx установите python-pptx]"
    
    prs = PptxPresentation(str(file_path))
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        texts.append(text)
        if texts:
            parts.append(f"--- Слайд {i} ---\n" + "\n".join(texts))
    return "\n\n".join(parts) if parts else "(Пустая презентация)"


def read_document(filename: str) -> str:
    print(f"[READ_DOC] Поиск файла: {filename}")
    
    # Ищем файл через общую функцию из file_operations
    file_path = find_file(filename)
    
    if not file_path:
        return f"Файл '{filename}' не найден. Проверьте имя файла."
    
    if not file_path.exists():
        return f"Файл '{filename}' не найден."
    
    print(f"[READ_DOC] Найден: {file_path}")
    
    # Определяем тип и читаем
    suffix = file_path.suffix.lower()
    
    try:
        if suffix in ('.txt', '.md', '.log', '.json', '.xml', '.html', '.css', '.js', '.py', '.csv'):
            content = _read_txt(file_path)
        elif suffix == '.docx':
            content = _read_docx(file_path)
        elif suffix == '.doc':
            content = _read_doc(file_path)
        elif suffix == '.pdf':
            content = _read_pdf(file_path)
        elif suffix in ('.xlsx', '.xls'):
            content = _read_xlsx(file_path)
        elif suffix == '.pptx':
            content = _read_pptx(file_path)
        else:
            # Пробуем как текст
            try:
                content = _read_txt(file_path)
            except Exception:
                return f"Формат файла '{suffix}' не поддерживается."
        
        # Обрезаем слишком длинный текст
        if len(content) > MAX_TEXT_LENGTH:
            content = content[:MAX_TEXT_LENGTH] + f"\n\n[... текст обрезан, всего {len(content)} символов]"
        
        if not content.strip():
            return f"Файл '{file_path.name}' пустой."
        
        return content
        
    except Exception as e:
        print(f"[READ_DOC] Ошибка чтения: {e}")
        return f"Ошибка чтения файла '{file_path.name}': {e}"


def execute_read_document(arguments: dict) -> str:
    filename = arguments.get("filename", "").strip()
    
    if not filename:
        return "Укажите имя файла для чтения."
    
    return read_document(filename)


def read_document_from_path(file_path: Path) -> str:
    """Читает документ напрямую по пути (для API загрузки файлов)."""
    if not file_path.exists():
        return f"Файл '{file_path.name}' не найден."
    
    suffix = file_path.suffix.lower()
    
    try:
        if suffix in ('.txt', '.md', '.log', '.json', '.xml', '.html', '.css', '.js', '.py', '.csv'):
            content = _read_txt(file_path)
        elif suffix == '.docx':
            content = _read_docx(file_path)
        elif suffix == '.doc':
            content = _read_doc(file_path)
        elif suffix == '.pdf':
            content = _read_pdf(file_path)
        elif suffix in ('.xlsx', '.xls'):
            content = _read_xlsx(file_path)
        elif suffix == '.pptx':
            content = _read_pptx(file_path)
        else:
            try:
                content = _read_txt(file_path)
            except Exception:
                return f"Формат файла '{suffix}' не поддерживается."
        
        if len(content) > MAX_TEXT_LENGTH:
            content = content[:MAX_TEXT_LENGTH] + f"\n\n[... текст обрезан, всего {len(content)} символов]"
        
        if not content.strip():
            return f"Файл '{file_path.name}' пустой."
        
        return content
        
    except Exception as e:
        print(f"[READ_DOC] Ошибка чтения: {e}")
        return f"Ошибка чтения файла '{file_path.name}': {e}"
