"""
Инструмент генерации документов для агента Вера.

Поддерживаемые форматы:
- .txt — текстовые файлы
- .md — Markdown документы
- .docx — Microsoft Word документы
- .pptx — PowerPoint презентации
- .xlsx — Excel таблицы

Все документы сохраняются в ~/Documents/Vera/
"""
import re
from pathlib import Path
from datetime import datetime
from typing import Dict, Any, List, Optional

# Опциональные библиотеки
try:
    from docx import Document as DocxDocument
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Inches, Pt
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

try:
    from openpyxl import Workbook
    from openpyxl.styles import Font, Alignment
    HAS_XLSX = True
except ImportError:
    HAS_XLSX = False


def get_documents_dir() -> Path:
    """
    Получает путь к папке для документов.
    Создаёт ~/Documents/Vera/ если не существует.
    """
    # Получаем путь к Documents
    try:
        # Windows
        import winreg
        with winreg.OpenKey(winreg.HKEY_CURRENT_USER, 
                           r"Software\Microsoft\Windows\CurrentVersion\Explorer\Shell Folders") as key:
            docs_path = Path(winreg.QueryValueEx(key, "Personal")[0])
    except Exception:
        # Fallback
        docs_path = Path.home() / "Documents"
    
    vera_docs = docs_path / "Vera"
    vera_docs.mkdir(parents=True, exist_ok=True)
    
    return vera_docs


def get_projects_dir() -> Path:
    """Возвращает папку для созданных презентаций."""
    projects_dir = get_documents_dir() / "Projects"
    projects_dir.mkdir(parents=True, exist_ok=True)
    return projects_dir


def sanitize_filename(name: str) -> str:
    """Очищает имя файла от недопустимых символов."""
    # Убираем расширение если есть
    name = re.sub(r'\.[a-zA-Z]{2,5}$', '', name)
    # Заменяем недопустимые символы
    name = re.sub(r'[<>:"/\\|?*]', '_', name)
    # Убираем лишние пробелы и подчёркивания
    name = re.sub(r'[\s_]+', '_', name).strip('_')
    # Ограничиваем длину
    if len(name) > 100:
        name = name[:100]
    return name or f"document_{datetime.now().strftime('%Y%m%d_%H%M%S')}"


def generate_unique_path(base_dir: Path, filename: str, extension: str) -> Path:
    """Генерирует уникальный путь к файлу (добавляет номер если файл существует)."""
    clean_name = sanitize_filename(filename)
    path = base_dir / f"{clean_name}.{extension}"
    
    if not path.exists():
        return path
    
    # Добавляем номер
    counter = 1
    while True:
        path = base_dir / f"{clean_name}_{counter}.{extension}"
        if not path.exists():
            return path
        counter += 1
        if counter > 100:  # Защита от бесконечного цикла
            path = base_dir / f"{clean_name}_{datetime.now().strftime('%H%M%S')}.{extension}"
            break
    
    return path


# ============ Генераторы документов ============

def create_txt(filename: str, content: str) -> str:
    """Создаёт текстовый файл."""
    docs_dir = get_documents_dir()
    file_path = generate_unique_path(docs_dir, filename, "txt")
    
    try:
        file_path.write_text(content, encoding='utf-8')
        return f"Файл создан: {file_path}"
    except Exception as e:
        return f"Ошибка создания файла: {e}"


def create_md(filename: str, content: str, title: Optional[str] = None) -> str:
    """Создаёт Markdown файл."""
    docs_dir = get_documents_dir()
    file_path = generate_unique_path(docs_dir, filename, "md")
    
    try:
        md_content = ""
        if title:
            md_content = f"# {title}\n\n"
        md_content += content
        
        file_path.write_text(md_content, encoding='utf-8')
        return f"Markdown создан: {file_path}"
    except Exception as e:
        return f"Ошибка создания Markdown: {e}"


def create_docx(
    filename: str,
    content: str,
    title: Optional[str] = None,
    author: Optional[str] = None
) -> str:
    """Создаёт Word документ."""
    if not HAS_DOCX:
        return "Ошибка: библиотека python-docx не установлена."
    
    docs_dir = get_documents_dir()
    file_path = generate_unique_path(docs_dir, filename, "docx")
    
    try:
        doc = DocxDocument()
        
        # Заголовок
        if title:
            heading = doc.add_heading(title, 0)
            heading.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        # Содержимое (разбиваем по абзацам)
        paragraphs = content.split('\n\n')
        for para in paragraphs:
            if para.strip():
                doc.add_paragraph(para.strip())
        
        # Метаданные
        if author:
            doc.core_properties.author = author
        doc.core_properties.created = datetime.now()
        
        doc.save(str(file_path))
        return f"Документ Word создан: {file_path}"
    except Exception as e:
        return f"Ошибка создания Word документа: {e}"


def create_pptx(
    filename: str,
    slides: List[Dict[str, Any]],
    title: Optional[str] = None,
    theme: str = "auto",
    subtitle: str = "",
    sources: Optional[List[str]] = None,
) -> str:
    """Create an editable 16:9 presentation with varied native layouts."""
    if not HAS_PPTX:
        return "Ошибка: библиотека python-pptx не установлена."

    docs_dir = get_projects_dir()
    file_path = generate_unique_path(docs_dir, filename, "pptx")

    palettes = {
        "ocean": {
            "bg": "F5F8FC", "panel": "FFFFFF", "ink": "122033",
            "muted": "617086", "accent": "177DDC", "accent2": "15A6A1",
            "soft": "E8F2FC", "family": "minimal",
            "title_font": "Aptos Display", "body_font": "Aptos",
        },
        "midnight": {
            "bg": "101724", "panel": "182337", "ink": "F7FAFF",
            "muted": "AAB7CA", "accent": "65C4FF", "accent2": "8B7CFF",
            "soft": "22324B", "family": "tech",
            "title_font": "Aptos Display", "body_font": "Aptos",
        },
        "forest": {
            "bg": "F5F7F2", "panel": "FFFFFF", "ink": "183126",
            "muted": "64766C", "accent": "23856D", "accent2": "D39B45",
            "soft": "E5F0EA", "family": "organic",
            "title_font": "Segoe UI Semibold", "body_font": "Segoe UI",
        },
        "violet": {
            "bg": "F8F6FC", "panel": "FFFFFF", "ink": "231C35",
            "muted": "746B84", "accent": "7657D6", "accent2": "D35493",
            "soft": "EEE9FA", "family": "bold",
            "title_font": "Arial", "body_font": "Aptos",
        },
        "sand": {
            "bg": "FBF7EF", "panel": "FFFDF8", "ink": "30271F",
            "muted": "796B5D", "accent": "D56B3F", "accent2": "287D8E",
            "soft": "F2E8D8", "family": "editorial",
            "title_font": "Georgia", "body_font": "Aptos",
        },
    }
    legacy_theme_map = {
        "white": "ocean", "light_blue": "ocean", "beige": "sand",
        "mint": "forest", "lavender": "violet",
    }
    if theme == "auto":
        topic_key = (title or filename).lower()
        content_key = " ".join(
            str(slide.get("title") or "") for slide in slides[:4]
        ).lower()
        variant = sum(ord(char) for char in content_key) % 2
        if any(word in topic_key for word in ("космос", "технолог", "ии", "ai", "кибер")):
            theme = ("midnight", "ocean")[variant]
        elif any(word in topic_key for word in ("природ", "эколог", "здоров", "био")):
            theme = ("forest", "sand")[variant]
        elif any(word in topic_key for word in ("искусств", "культур", "дизайн", "мода")):
            theme = ("violet", "sand")[variant]
        elif any(word in topic_key for word in ("истор", "литератур", "образован")):
            theme = ("sand", "violet")[variant]
        else:
            theme = ("ocean", "forest", "sand", "violet")[sum(ord(char) for char in topic_key) % 4]
    theme = legacy_theme_map.get(theme, theme)
    colors = palettes.get(theme, palettes["ocean"])

    def rgb(key: str) -> RGBColor:
        value = colors.get(key, key).lstrip("#")
        return RGBColor.from_string(value)

    def add_text(
        slide,
        text: str,
        x: float,
        y: float,
        w: float,
        h: float,
        size: float,
        color: str = "ink",
        bold: bool = False,
        font: Optional[str] = None,
        align=PP_ALIGN.LEFT,
        valign=MSO_ANCHOR.TOP,
    ):
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        frame = box.text_frame
        frame.clear()
        frame.word_wrap = True
        frame.margin_left = frame.margin_right = Inches(0.02)
        frame.margin_top = frame.margin_bottom = Inches(0.02)
        frame.vertical_anchor = valign
        paragraph = frame.paragraphs[0]
        paragraph.alignment = align
        run = paragraph.add_run()
        run.text = str(text or "")
        run.font.name = font or (
            colors["title_font"] if bold and size >= 20 else colors["body_font"]
        )
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = rgb(color)
        return box

    def add_rect(
        slide,
        x: float,
        y: float,
        w: float,
        h: float,
        fill: str = "panel",
        radius: Optional[bool] = None,
        line: Optional[str] = None,
    ):
        if radius is None:
            radius = colors["family"] in ("tech", "organic")
        shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
        shape = slide.shapes.add_shape(
            shape_type, Inches(x), Inches(y), Inches(w), Inches(h)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(fill)
        shape.line.color.rgb = rgb(line or fill)
        shape.line.width = Pt(0.8)
        return shape

    def add_circle(slide, x: float, y: float, d: float, fill: str = "accent"):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(fill)
        shape.line.color.rgb = rgb(fill)
        return shape

    def add_base(slide, index: int, kicker: str = ""):
        background = slide.background.fill
        background.solid()
        background.fore_color.rgb = rgb("bg")
        family = colors["family"]
        if family == "tech":
            add_rect(slide, 0, 0, 0.11, 7.5, "accent", False, "accent")
            add_circle(slide, 11.85, -0.55, 1.45, "soft")
        elif family == "editorial":
            add_rect(slide, 0.65, 0.45, 12.0, 0.025, "accent", False, "accent")
        elif family == "bold":
            add_rect(slide, 0, 0, 2.15, 0.18, "accent2", False, "accent2")
            add_rect(slide, 2.15, 0, 3.0, 0.18, "accent", False, "accent")
        elif family == "organic":
            add_circle(slide, 11.75, -0.65, 1.7, "soft")
            add_circle(slide, 12.45, 0.1, 0.42, "accent2")
        else:
            add_rect(slide, 0.65, 0.45, 0.8, 0.05, "accent", False, "accent")
        if kicker:
            kicker_x = 0.72 if family != "editorial" else 9.0
            kicker_w = 4.5 if family != "editorial" else 3.6
            add_text(
                slide, kicker.upper(), kicker_x, 0.38, kicker_w, 0.25,
                9, "accent", True,
                align=PP_ALIGN.RIGHT if family == "editorial" else PP_ALIGN.LEFT,
            )
        add_text(
            slide, f"{index:02d}", 12.25, 7.05, 0.45, 0.2,
            8, "muted", True, align=PP_ALIGN.RIGHT,
        )

    def add_heading(slide, data: Dict[str, Any], index: int):
        add_base(slide, index, str(data.get("kicker") or ""))
        add_text(slide, data.get("title", ""), 0.62, 0.72, 11.6, 0.72, 26, "ink", True)

    def clean_bullets(data: Dict[str, Any]) -> List[str]:
        bullets = data.get("bullets") or []
        if isinstance(bullets, str):
            bullets = [part.strip() for part in bullets.split("\n") if part.strip()]
        return [str(item).strip() for item in bullets if str(item).strip()][:4]

    def render_overview(slide, data: Dict[str, Any], index: int):
        add_heading(slide, data, index)
        message = str(data.get("key_message") or data.get("content") or "")
        add_text(slide, message, 0.65, 1.65, 11.3, 0.78, 18, "ink", True)
        bullets = clean_bullets(data) or ["Контекст", "Ключевая идея", "Практический вывод"]
        if colors["family"] == "editorial":
            for item_index, item in enumerate(bullets):
                y = 2.72 + item_index * 0.86
                add_text(slide, f"0{item_index + 1}", 0.72, y, 0.55, 0.35, 11, "accent", True)
                add_text(slide, item, 1.55, y - 0.05, 9.8, 0.55, 16, "ink", item_index == 0)
                add_rect(slide, 1.55, y + 0.58, 10.5, 0.015, "soft", False, "soft")
            return
        if colors["family"] == "bold":
            for item_index, item in enumerate(bullets):
                x = 0.65 + (item_index % 2) * 5.95
                y = 2.65 + (item_index // 2) * 1.55
                fill = "accent" if item_index % 2 == 0 else "accent2"
                add_rect(slide, x, y, 5.55, 1.2, fill, False, fill)
                add_text(slide, item, x + 0.28, y + 0.32, 4.95, 0.55, 15, "panel", True)
            return
        count = len(bullets)
        card_w = (11.8 - (count - 1) * 0.28) / count
        for item_index, item in enumerate(bullets):
            x = 0.65 + item_index * (card_w + 0.28)
            add_rect(slide, x, 2.75, card_w, 2.75, "panel", True, "soft")
            add_circle(slide, x + 0.24, 3.02, 0.48, "accent" if item_index % 2 == 0 else "accent2")
            add_text(slide, str(item_index + 1), x + 0.24, 3.08, 0.48, 0.22, 10, "panel", True, align=PP_ALIGN.CENTER)
            add_text(slide, item, x + 0.25, 3.72, card_w - 0.5, 1.35, 15, "ink", True)

    def render_process(slide, data: Dict[str, Any], index: int, timeline: bool = False):
        add_heading(slide, data, index)
        add_text(slide, data.get("key_message", ""), 0.65, 1.58, 11.3, 0.58, 16, "muted")
        items = clean_bullets(data) or ["Исходная точка", "Ключевой шаг", "Результат"]
        count = len(items)
        y = 3.65
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(1.2), Inches(y + 0.24), Inches(10.5), Inches(0.04)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = rgb("soft")
        line.line.fill.background()
        for item_index, item in enumerate(items):
            x = 0.85 + item_index * (10.9 / max(1, count - 1)) if count > 1 else 6.1
            add_circle(slide, x, y, 0.52, "accent" if item_index % 2 == 0 else "accent2")
            add_text(slide, str(item_index + 1), x, y + 0.12, 0.52, 0.2, 10, "panel", True, align=PP_ALIGN.CENTER)
            label_y = 2.55 if timeline and item_index % 2 == 0 else 4.45
            add_text(slide, item, x - 0.7, label_y, 1.95, 0.85, 12, "ink", True, align=PP_ALIGN.CENTER)

    def render_comparison(slide, data: Dict[str, Any], index: int):
        add_heading(slide, data, index)
        add_text(slide, data.get("key_message", ""), 0.65, 1.55, 11.3, 0.6, 16, "muted")
        bullets = clean_bullets(data)
        midpoint = max(1, (len(bullets) + 1) // 2)
        columns = [bullets[:midpoint], bullets[midpoint:]]
        labels = ["С одной стороны", "С другой стороны"]
        fills = ["soft", "panel"]
        for col in range(2):
            x = 0.65 + col * 6.05
            add_rect(slide, x, 2.45, 5.7, 3.65, fills[col], True, "soft")
            add_text(slide, labels[col], x + 0.35, 2.8, 4.9, 0.4, 12, "accent" if col == 0 else "accent2", True)
            for row, item in enumerate(columns[col] or ["Добавьте критерий сравнения"]):
                add_circle(slide, x + 0.38, 3.55 + row * 0.78, 0.16, "accent" if col == 0 else "accent2")
                add_text(slide, item, x + 0.72, 3.42 + row * 0.78, 4.45, 0.55, 13, "ink", row == 0)

    def render_numbers(slide, data: Dict[str, Any], index: int):
        add_heading(slide, data, index)
        add_text(slide, data.get("key_message", ""), 0.65, 1.55, 11.3, 0.6, 16, "muted")
        stats = data.get("stats") if isinstance(data.get("stats"), list) else []
        if not stats:
            stats = [
                {"value": f"{i + 1:02d}", "label": item}
                for i, item in enumerate(clean_bullets(data)[:3])
            ]
        stats = stats[:3] or [{"value": "01", "label": "Главный ориентир"}]
        card_w = 3.55
        start_x = (13.33 - (card_w * len(stats) + 0.35 * (len(stats) - 1))) / 2
        for item_index, stat in enumerate(stats):
            x = start_x + item_index * (card_w + 0.35)
            add_rect(slide, x, 2.55, card_w, 3.15, "panel", True, "soft")
            add_text(slide, stat.get("value", ""), x + 0.25, 3.0, card_w - 0.5, 0.85, 30, "accent" if item_index != 1 else "accent2", True, align=PP_ALIGN.CENTER)
            add_text(slide, stat.get("label", ""), x + 0.35, 4.15, card_w - 0.7, 0.85, 13, "ink", True, align=PP_ALIGN.CENTER)

    def render_quote(slide, data: Dict[str, Any], index: int):
        add_heading(slide, data, index)
        quote = data.get("quote") or data.get("key_message") or ""
        add_text(slide, "“", 0.75, 1.75, 0.8, 0.8, 54, "accent", True)
        add_text(slide, quote, 1.55, 2.05, 9.9, 2.1, 26, "ink", True, valign=MSO_ANCHOR.MIDDLE)
        bullets = clean_bullets(data)
        if bullets:
            add_text(slide, " • ".join(bullets[:3]), 1.58, 4.65, 9.8, 0.65, 12, "muted")

    def render_summary(slide, data: Dict[str, Any], index: int):
        add_heading(slide, data, index)
        add_rect(slide, 0.65, 1.65, 12.0, 1.55, "accent", True, "accent")
        add_text(slide, data.get("key_message", ""), 1.0, 1.95, 11.3, 0.9, 22, "panel", True, valign=MSO_ANCHOR.MIDDLE)
        for row, item in enumerate(clean_bullets(data)[:3]):
            y = 3.65 + row * 0.78
            add_circle(slide, 0.85, y, 0.36, "accent2")
            add_text(slide, str(row + 1), 0.85, y + 0.07, 0.36, 0.17, 9, "panel", True, align=PP_ALIGN.CENTER)
            add_text(slide, item, 1.45, y - 0.03, 10.5, 0.5, 15, "ink", True)

    try:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        blank = prs.slide_layouts[6]

        if title:
            slide = prs.slides.add_slide(blank)
            background = slide.background.fill
            background.solid()
            background.fore_color.rgb = rgb("bg")
            family = colors["family"]
            if family == "tech":
                add_circle(slide, 9.4, -1.0, 4.8, "soft")
                add_circle(slide, 10.55, 0.15, 2.65, "accent")
            elif family == "editorial":
                add_rect(slide, 9.85, 0, 3.48, 7.5, "soft", False, "soft")
                add_rect(slide, 0.78, 1.2, 1.25, 0.045, "accent", False, "accent")
            elif family == "bold":
                add_rect(slide, 8.75, 0, 4.58, 7.5, "accent", False, "accent")
                add_rect(slide, 9.55, 0.72, 3.15, 1.45, "accent2", False, "accent2")
                add_text(slide, "01", 9.78, 3.0, 2.5, 1.4, 54, "panel", True)
            elif family == "organic":
                add_circle(slide, 9.05, -0.75, 5.2, "soft")
                add_circle(slide, 10.4, 0.65, 2.75, "accent2")
                add_circle(slide, 11.35, 1.6, 1.35, "accent")
            else:
                add_rect(slide, 0.78, 1.35, 0.1, 3.25, "accent", False, "accent")
            add_text(slide, "VERA • PRESENTATION", 0.75, 0.65, 4.0, 0.3, 10, "accent", True)
            title_width = 8.35 if family in ("editorial", "bold") else 8.9
            title_size = 37 if family == "editorial" else (40 if family == "bold" else 34)
            add_text(slide, title, 0.75, 1.55, title_width, 2.3, title_size, "ink", True, valign=MSO_ANCHOR.MIDDLE)
            add_text(
                slide,
                subtitle or "Ключевые идеи, факты и практические выводы",
                0.78,
                4.25,
                7.8,
                0.75,
                16,
                "muted",
            )
            add_rect(slide, 0.78, 5.45, 1.2, 0.08, "accent", False, "accent")
            add_text(slide, datetime.now().strftime("%d.%m.%Y"), 0.78, 5.75, 2.0, 0.3, 10, "muted")

        renderers = {
            "overview": render_overview,
            "process": render_process,
            "comparison": render_comparison,
            "numbers": render_numbers,
            "timeline": lambda slide, data, index: render_process(slide, data, index, True),
            "quote": render_quote,
            "summary": render_summary,
        }
        for index, slide_data in enumerate(slides, start=2 if title else 1):
            slide = prs.slides.add_slide(blank)
            visual = str(slide_data.get("visual") or "overview").lower()
            renderer = renderers.get(visual, render_overview)
            renderer(slide, slide_data, index)

        if sources and prs.slides:
            from urllib.parse import urlparse
            domains = []
            for source in sources:
                domain = urlparse(source).netloc.removeprefix("www.")
                if domain and domain not in domains:
                    domains.append(domain)
            if domains:
                add_text(
                    prs.slides[-1],
                    "Источники: " + " • ".join(domains[:5]),
                    0.65,
                    7.05,
                    10.8,
                    0.18,
                    7,
                    "muted",
                )

        prs.save(str(file_path))
        return f"Презентация создана: {file_path} (стиль: {theme})"
    except Exception as e:
        return f"Ошибка создания презентации: {e}"


def create_xlsx(
    filename: str,
    data: List[List[Any]],
    headers: Optional[List[str]] = None,
    sheet_name: str = "Лист1"
) -> str:
    """
    Создаёт Excel таблицу.
    
    Args:
        filename: Имя файла
        data: Данные [[row1], [row2], ...]
        headers: Заголовки столбцов
        sheet_name: Название листа
    """
    if not HAS_XLSX:
        return "Ошибка: библиотека openpyxl не установлена."
    
    docs_dir = get_documents_dir()
    file_path = generate_unique_path(docs_dir, filename, "xlsx")
    
    try:
        wb = Workbook()
        ws = wb.active
        ws.title = sheet_name
        
        start_row = 1
        
        # Заголовки
        if headers:
            for col, header in enumerate(headers, 1):
                cell = ws.cell(row=1, column=col, value=header)
                cell.font = Font(bold=True)
                cell.alignment = Alignment(horizontal='center')
            start_row = 2
        
        # Данные
        for row_idx, row_data in enumerate(data, start_row):
            for col_idx, value in enumerate(row_data, 1):
                ws.cell(row=row_idx, column=col_idx, value=value)
        
        # Автоширина столбцов
        for column_cells in ws.columns:
            length = max(len(str(cell.value or "")) for cell in column_cells)
            ws.column_dimensions[column_cells[0].column_letter].width = min(length + 2, 50)
        
        wb.save(str(file_path))
        return f"Таблица Excel создана: {file_path}"
    except Exception as e:
        return f"Ошибка создания Excel таблицы: {e}"


# ============ Диспетчер инструмента ============

def execute_document_generator(arguments: dict) -> str:
    """
    Главная функция инструмента — диспетчер по действиям.
    
    Поддерживаемые действия:
    - create_txt: {"action": "create_txt", "filename": "...", "content": "..."}
    - create_md: {"action": "create_md", "filename": "...", "content": "...", "title": "..."}
    - create_docx: {"action": "create_docx", "filename": "...", "content": "...", "title": "..."}
    - create_pptx: {"action": "create_pptx", "filename": "...", "slides": [...], "title": "..."}
    - create_xlsx: {"action": "create_xlsx", "filename": "...", "data": [...], "headers": [...]}
    """
    action = arguments.get("action", "").strip().lower()
    filename = arguments.get("filename", "").strip()
    
    if not action:
        return "Укажите действие: create_txt, create_md, create_docx, create_pptx, create_xlsx"
    
    if not filename:
        filename = f"document_{datetime.now().strftime('%Y%m%d_%H%M%S')}"
    
    print(f"[DOC_GEN] Действие: {action}, файл: {filename}")
    
    # Текстовый файл
    if action == "create_txt":
        content = arguments.get("content", "")
        if not content:
            return "Укажите content — текст для файла."
        return create_txt(filename, content)
    
    # Markdown
    if action == "create_md":
        content = arguments.get("content", "")
        title = arguments.get("title")
        if not content:
            return "Укажите content — текст для Markdown."
        return create_md(filename, content, title)
    
    # Word документ
    if action == "create_docx":
        content = arguments.get("content", "")
        title = arguments.get("title")
        author = arguments.get("author")
        if not content:
            return "Укажите content — текст документа."
        return create_docx(filename, content, title, author)
    
    # PowerPoint
    if action == "create_pptx":
        slides = arguments.get("slides", [])
        title = arguments.get("title")
        # Если вместо slides передан content — преобразуем
        if not slides and "content" in arguments:
            content = arguments["content"]
            slides = [{"title": title or "Слайд 1", "content": content}]
        if not slides:
            return "Укажите slides — список слайдов."
        return create_pptx(filename, slides, title, arguments.get("theme", "light_blue"))
    
    # Excel
    if action == "create_xlsx":
        data = arguments.get("data", [])
        headers = arguments.get("headers")
        sheet_name = arguments.get("sheet_name", "Лист1")
        if not data:
            return "Укажите data — данные таблицы."
        return create_xlsx(filename, data, headers, sheet_name)
    
    return f"Неизвестное действие: {action}. Доступны: create_txt, create_md, create_docx, create_pptx, create_xlsx"
